#!/usr/bin/env python3
"""
NUS Canvas Deadline Scraper
Automatically extracts assignment deadlines and exam dates from Canvas course materials.
"""

import os
import json
import logging
import argparse
import re
import time
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Any, Optional
import requests
from dotenv import load_dotenv
import google.generativeai as genai
from html.parser import HTMLParser

# Import document parsers
try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False
    logging.warning("PyMuPDF not available - PDF parsing disabled")

try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False
    logging.warning("pdfplumber not available - advanced PDF parsing disabled")

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logging.warning("python-docx not available - DOCX parsing disabled")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not available - PPTX parsing disabled")

# Setup logging (console only, no log file clutter)
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

# Load environment variables
load_dotenv()

# Configuration
CANVAS_API_TOKEN = os.getenv('CANVAS_API_TOKEN')
CANVAS_API_URL = os.getenv('CANVAS_API_URL', 'https://canvas.nus.edu.sg/api/v1')
GEMINI_API_KEY = os.getenv('GEMINI_API_KEY')

# Auto-detect current semester based on today's date
def get_current_semester_bounds():
    """
    Auto-detect current NUS semester based on current date.
    NUS has 2 main semesters:
    - Semester 1: August - December
    - Semester 2: January - May
    - Special Term: June - July (optional)
    """
    today = datetime.now()
    year = today.year
    month = today.month

    # Allow manual override via .env
    manual_start = os.getenv('SEMESTER_START')
    manual_end = os.getenv('SEMESTER_END')

    if manual_start and manual_end:
        return manual_start, manual_end

    # Auto-detect based on current month
    if 1 <= month <= 5:
        # Semester 2: January - May
        return f"{year}-01-01", f"{year}-05-31"
    elif 8 <= month <= 12:
        # Semester 1: August - December
        return f"{year}-08-01", f"{year}-12-31"
    else:
        # June-July: Special Term or between semesters
        # Default to upcoming Semester 1
        return f"{year}-08-01", f"{year}-12-31"

SEMESTER_START, SEMESTER_END = get_current_semester_bounds()

# Keywords to identify intro/syllabus documents
INTRO_KEYWORDS = [
    # Basic intro keywords
    'intro', 'introduction', 'syllabus', 'course outline',
    'schedule', 'overview', 'course info', 'course information',
    'module information', 'module info',

    # Schedule and assessment keywords
    'course schedule', 'assessment', 'assessments',

    # Week/Lecture variations
    'week 1', 'week1', 'week 0', 'week0',
    'lecture 1', 'lecture1', 'lecture 0', 'lecture0',
    'lec 1', 'lec1', 'lec 0', 'lec0',

    # Topic variations
    'topic 0', 'topic_0', 'topic0',

    # Year/semester indicators combined with module codes
    'ay24', 'ay25', 'ay26', 's1', 's2',
]

# Regex patterns for intro documents
INTRO_PATTERNS = [
    # Lecture number patterns: L0, L00, L01, L1, L01a, etc.
    r'[_\s\-]?l0+[a-z]?[_\s\-]',  # L0, L00, L0a, L00a
    r'[_\s\-]?l0*1[a-z]?[_\s\-]',  # L1, L01, L001, L1a, L01a

    # Lec/Lecture patterns
    r'lec[_\s\-]?0+[_\s\-]',  # lec0, lec00, lec_0
    r'lec[_\s\-]?0*1[_\s\-]',  # lec1, lec01, lec_1

    # Topic patterns
    r'topic[_\s\-]?0',  # topic0, topic_0, topic 0

    # Week patterns
    r'week[_\s\-]?0*1[_\s\-]',  # week1, week01, week_1
]

# Directory for downloaded files (temporary cache)
DOWNLOAD_DIR = Path('.temp')
DOWNLOAD_DIR.mkdir(exist_ok=True)


class HTMLTextExtractor(HTMLParser):
    """Extract plain text from HTML, preserving table structure"""
    def __init__(self):
        super().__init__()
        self.text = []
        self.in_table = False
        self.in_row = False
        self.in_cell = False
        self.current_row = []

    def handle_starttag(self, tag, attrs):
        if tag == 'table':
            self.in_table = True
            self.text.append('\n\nTABLE:\n')
        elif tag == 'tr' and self.in_table:
            self.in_row = True
            self.current_row = []
        elif tag in ['td', 'th'] and self.in_row:
            self.in_cell = True
        elif tag == 'br':
            self.text.append('\n')

    def handle_endtag(self, tag):
        if tag == 'table':
            self.in_table = False
            self.text.append('\nEND TABLE\n\n')
        elif tag == 'tr' and self.in_row:
            self.in_row = False
            if self.current_row:
                self.text.append(' | '.join(self.current_row) + '\n')
            self.current_row = []
        elif tag in ['td', 'th'] and self.in_cell:
            self.in_cell = False
        elif tag == 'p':
            self.text.append('\n')

    def handle_data(self, data):
        data = data.strip()
        if data:
            if self.in_cell:
                self.current_row.append(data)
            else:
                self.text.append(data + ' ')

    def get_text(self):
        return ''.join(self.text)


def html_to_text(html: str) -> str:
    """Convert HTML to plain text, preserving table structure"""
    parser = HTMLTextExtractor()
    parser.feed(html)
    return parser.get_text()


def is_valid_semester_date(date_str: str) -> bool:
    """
    Validate if a date falls within the configured semester bounds.
    Filters out deadlines from previous semesters that might be in old materials.
    """
    if not date_str or date_str == 'TBD':
        return True  # Allow TBD dates

    try:
        deadline_date = datetime.fromisoformat(date_str).date()
        semester_start = datetime.fromisoformat(SEMESTER_START).date()
        semester_end = datetime.fromisoformat(SEMESTER_END).date()

        return semester_start <= deadline_date <= semester_end
    except (ValueError, TypeError):
        # If date parsing fails, keep the deadline (be permissive)
        logger.warning(f"Could not validate date format: {date_str}")
        return True


def filter_valid_semester_deadlines(deadlines: List[Dict]) -> tuple[List[Dict], List[Dict]]:
    """
    Filter deadlines to only include those within the semester bounds.
    Returns (valid_deadlines, filtered_out_deadlines)
    """
    valid = []
    filtered_out = []

    for deadline in deadlines:
        date_str = deadline.get('date', '')
        if is_valid_semester_date(date_str):
            valid.append(deadline)
        else:
            filtered_out.append(deadline)
            logger.debug(f"Filtered out deadline outside semester: {deadline.get('title')} ({date_str})")

    return valid, filtered_out


class CanvasClient:
    """Client for Canvas LMS API"""

    def __init__(self, api_url: str, api_token: str):
        self.api_url = api_url.rstrip('/')
        self.headers = {
            'Authorization': f'Bearer {api_token}',
            'Accept': 'application/json'
        }

    def _get(self, endpoint: str, params: Optional[Dict] = None) -> Any:
        """Make GET request to Canvas API"""
        url = f"{self.api_url}/{endpoint.lstrip('/')}"
        try:
            response = requests.get(url, headers=self.headers, params=params)
            response.raise_for_status()
            return response.json()
        except requests.exceptions.RequestException as e:
            logger.error(f"API request failed: {e}")
            return None

    def get_current_user(self) -> Dict:
        """Get current user information"""
        return self._get('/users/self')

    def get_enrolled_courses(self) -> List[Dict]:
        """Get all enrolled courses for current user"""
        courses = self._get('/courses', params={
            'enrollment_state': 'active',
            'per_page': 100
        })
        return courses if courses else []

    def get_course_files(self, course_id: int) -> List[Dict]:
        """Get all files for a course"""
        files = self._get(f'/courses/{course_id}/files', params={
            'per_page': 100
        })
        return files if files else []

    def get_course_assignments(self, course_id: int) -> List[Dict]:
        """Get all assignments for a course"""
        assignments = self._get(f'/courses/{course_id}/assignments', params={
            'per_page': 100
        })
        return assignments if assignments else []

    def get_course_quizzes(self, course_id: int) -> List[Dict]:
        """Get all quizzes for a course"""
        quizzes = self._get(f'/courses/{course_id}/quizzes', params={
            'per_page': 100
        })
        return quizzes if quizzes else []

    def get_course_details(self, course_id: int) -> Dict:
        """Get course details including syllabus/homepage"""
        course = self._get(f'/courses/{course_id}', params={
            'include[]': 'syllabus_body'
        })
        return course if course else {}

    def get_course_announcements(self, course_id: int, limit: int = 50) -> List[Dict]:
        """Get recent announcements for a course"""
        announcements = self._get(f'/courses/{course_id}/discussion_topics', params={
            'only_announcements': True,
            'per_page': limit,
            'order_by': 'recent_activity'
        })
        return announcements if announcements else []

    def get_calendar_events(self, course_id: Optional[int] = None) -> List[Dict]:
        """Get calendar events (all courses or specific course)"""
        params = {
            'per_page': 100,
            'type': 'event'
        }
        if course_id:
            params['context_codes[]'] = f'course_{course_id}'

        events = self._get('/calendar_events', params=params)
        return events if events else []

    def get_course_modules(self, course_id: int) -> List[Dict]:
        """Get all modules for a course"""
        modules = self._get(f'/courses/{course_id}/modules', params={
            'per_page': 100,
            'include[]': 'items'
        })
        return modules if modules else []

    def download_file(self, file_url: str, save_path: Path) -> bool:
        """Download a file from Canvas"""
        try:
            response = requests.get(file_url, headers=self.headers, stream=True)
            response.raise_for_status()

            with open(save_path, 'wb') as f:
                for chunk in response.iter_content(chunk_size=8192):
                    f.write(chunk)

            logger.info(f"Downloaded: {save_path.name}")
            return True
        except Exception as e:
            logger.error(f"Failed to download {save_path.name}: {e}")
            return False


class DocumentParser:
    """Parse different document formats to extract text"""

    @staticmethod
    def parse_pdf(file_path: Path) -> str:
        """Extract text from PDF using PyMuPDF (better than PyPDF2)"""
        if not PYMUPDF_AVAILABLE:
            return ""

        try:
            # Use PyMuPDF for better text extraction
            doc = fitz.open(str(file_path))
            text = ""
            for page in doc:
                text += page.get_text() + "\n"
            doc.close()
            return text
        except Exception as e:
            logger.error(f"Failed to parse PDF {file_path.name}: {e}")
            return ""

    @staticmethod
    def parse_docx(file_path: Path) -> str:
        """Extract text from DOCX"""
        if not DOCX_AVAILABLE:
            return ""

        try:
            doc = DocxDocument(str(file_path))
            return "\n".join([para.text for para in doc.paragraphs])
        except Exception as e:
            logger.error(f"Failed to parse DOCX {file_path.name}: {e}")
            return ""

    @staticmethod
    def parse_pptx(file_path: Path) -> str:
        """Extract text from PowerPoint"""
        if not PPTX_AVAILABLE:
            return ""

        try:
            prs = Presentation(str(file_path))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Failed to parse PPTX {file_path.name}: {e}")
            return ""

    @staticmethod
    def parse_txt(file_path: Path) -> str:
        """Extract text from TXT file"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception as e:
            logger.error(f"Failed to parse TXT {file_path.name}: {e}")
            return ""

    @classmethod
    def parse_document(cls, file_path: Path) -> str:
        """Parse document based on file extension"""
        ext = file_path.suffix.lower()

        parsers = {
            '.pdf': cls.parse_pdf,
            '.docx': cls.parse_docx,
            '.doc': cls.parse_docx,
            '.pptx': cls.parse_pptx,
            '.ppt': cls.parse_pptx,
            '.txt': cls.parse_txt,
        }

        parser = parsers.get(ext)
        if parser:
            return parser(file_path)
        else:
            logger.warning(f"Unsupported file format: {ext}")
            return ""


class DeadlineExtractor:
    """Extract deadlines from documents using Gemini AI with automatic model rotation"""

    # Available Gemini models (each has separate FREE quota!)
    # Using "latest" aliases ensures auto-updates when Google releases new versions
    AVAILABLE_MODELS = [
        'gemini-2.5-flash',                    # Stable: balanced performance
        'gemini-flash-latest',                 # Auto-updates to latest Flash (hot-swapped by Google)
        'gemini-2.5-flash-lite',               # Stable: fastest, most cost-effective
        'gemini-2.5-pro',                      # Stable: most powerful
        'gemini-pro-latest',                   # Auto-updates to latest Pro
    ]

    def __init__(self, api_key: str):
        genai.configure(api_key=api_key)
        self.api_key = api_key
        self.current_model_index = 0
        self.available_models = self.AVAILABLE_MODELS.copy()

        # Try to fetch latest models dynamically as fallback
        try:
            self._fetch_available_models()
        except Exception as e:
            logger.warning(f"Could not fetch dynamic model list, using hardcoded: {e}")

        self.model = genai.GenerativeModel(self.available_models[self.current_model_index])
        logger.info(f"Initialized with model: {self.available_models[self.current_model_index]}")
        logger.info(f"Total models available: {len(self.available_models)}")

    def _fetch_available_models(self):
        """Dynamically fetch available Gemini models from API"""
        try:
            # List available models from Gemini API
            models = genai.list_models()

            # Filter for models that support generateContent
            gemini_models = [
                m.name.replace('models/', '')
                for m in models
                if 'generateContent' in m.supported_generation_methods
                and 'gemini' in m.name.lower()
            ]

            if gemini_models:
                # Add newly discovered models to the list (avoid duplicates)
                for model in gemini_models:
                    if model not in self.available_models:
                        self.available_models.append(model)
                        logger.info(f"Discovered new model: {model}")
        except Exception as e:
            logger.debug(f"Model discovery failed: {e}")

    def _generate_with_fallback(self, prompt: str, generation_config: dict, max_retries: int = 3):
        """Generate content with automatic model rotation on quota exceeded"""
        attempts = 0
        last_error = None

        while attempts < max_retries:
            try:
                # Add delay to respect rate limits
                time.sleep(6.5)

                response = self.model.generate_content(prompt, generation_config=generation_config)
                return response

            except Exception as e:
                error_str = str(e)

                # Check if it's a quota error (429)
                if '429' in error_str and 'quota' in error_str.lower():
                    logger.warning(f"Quota exceeded for {self.available_models[self.current_model_index]}")

                    # Try next model
                    self.current_model_index = (self.current_model_index + 1) % len(self.available_models)
                    new_model_name = self.available_models[self.current_model_index]

                    logger.info(f"🔄 Switching to model: {new_model_name}")
                    self.model = genai.GenerativeModel(new_model_name)

                    attempts += 1
                    last_error = e

                    # If we've tried all models, wait a bit before retrying
                    if attempts >= len(self.available_models):
                        logger.warning(f"All models exhausted, waiting 10 seconds before retry {attempts}/{max_retries}...")
                        time.sleep(10)
                else:
                    # Not a quota error, raise immediately
                    raise e

        # If we exhausted all retries, raise the last error
        if last_error:
            raise last_error

    def extract_deadlines(self, text: str, course_name: str, current_year: int) -> List[Dict]:
        """Use Gemini to extract deadlines from document text"""

        if not text.strip():
            logger.warning(f"No text to extract from {course_name}")
            return []

        # Truncate text if too long (Gemini has token limits)
        max_chars = 100000
        if len(text) > max_chars:
            text = text[:max_chars] + "\n... [truncated]"

        prompt = f"""You are analyzing a course document for "{course_name}".
Extract ALL assignment deadlines, exam dates, quiz dates, project deadlines, and any other important academic dates.

IMPORTANT - Current semester context:
- Semester runs from {SEMESTER_START} to {SEMESTER_END}
- Current year: {current_year}
- ONLY extract deadlines that fall within the current semester period
- IGNORE any dates from previous semesters (e.g., dates before {SEMESTER_START})
- If a document contains both old and new deadlines, ONLY extract the current semester ones

Document text:
{text}

Please extract and return ONLY a JSON array of deadline objects. Each object should have:
- "date": ISO format date (YYYY-MM-DD) - MUST be between {SEMESTER_START} and {SEMESTER_END}
- "title": Brief description of the deadline
- "type": One of ["assignment", "exam", "quiz", "project", "presentation", "other"]
- "weight": Percentage weight if mentioned (or null)
- "notes": Any additional relevant information

Return ONLY the JSON array, no other text. If no deadlines found, return an empty array [].

Example format:
[
  {{"date": "2025-09-15", "title": "Assignment 1", "type": "assignment", "weight": 10, "notes": "Submit via Canvas"}},
  {{"date": "2025-11-20", "title": "Final Exam", "type": "exam", "weight": 50, "notes": "Open book"}}
]
"""

        try:
            logger.info(f"Extracting deadlines from {course_name} using Gemini...")

            # Use fallback system for automatic model rotation
            response = self._generate_with_fallback(
                prompt,
                generation_config=genai.types.GenerationConfig(
                    temperature=0,
                    max_output_tokens=4096,
                )
            )

            response_text = response.text.strip()

            # Try to parse JSON response
            # Remove markdown code blocks if present
            if response_text.startswith('```'):
                response_text = response_text.split('```')[1]
                if response_text.startswith('json'):
                    response_text = response_text[4:]
                response_text = response_text.strip()

            deadlines = json.loads(response_text)

            # Add course name to each deadline
            for deadline in deadlines:
                deadline['course'] = course_name

            logger.info(f"Extracted {len(deadlines)} deadlines from {course_name}")
            return deadlines

        except json.JSONDecodeError as e:
            logger.error(f"Failed to parse JSON response for {course_name}: {e}")
            logger.debug(f"Response was: {response_text}")
            return []
        except Exception as e:
            logger.error(f"Failed to extract deadlines for {course_name}: {e}")
            return []

    def extract_from_announcements(self, announcements: List[Dict], course_name: str, current_year: int) -> List[Dict]:
        """Extract deadline updates from Canvas announcements"""
        all_updates = []

        for announcement in announcements:
            title = announcement.get('title', '')
            message = announcement.get('message', '')
            posted_at = announcement.get('posted_at', '')

            # Combine title and message
            full_text = f"Title: {title}\n\n{message}"

            # Skip if no deadline-related keywords
            deadline_keywords = ['deadline', 'due', 'extended', 'postponed', 'rescheduled',
                                'changed', 'updated', 'new date', 'assignment', 'exam', 'quiz']
            if not any(keyword in full_text.lower() for keyword in deadline_keywords):
                continue

            logger.info(f"  Checking announcement: {title}")

            # Use AI to extract deadline updates
            prompt = f"""You are analyzing a Canvas announcement from "{course_name}" posted on {posted_at}.
Extract any deadline changes, updates, or new deadlines mentioned.

IMPORTANT - Current semester context:
- Semester runs from {SEMESTER_START} to {SEMESTER_END}
- Current year: {current_year}
- ONLY extract deadlines within the current semester period
- IGNORE any old dates from previous semesters

Announcement:
{full_text[:5000]}

Please extract and return ONLY a JSON array of deadline objects. Each object should have:
- "date": ISO format date (YYYY-MM-DD) - MUST be between {SEMESTER_START} and {SEMESTER_END}
- "title": Brief description
- "type": One of ["assignment", "exam", "quiz", "project", "presentation", "other"]
- "weight": Percentage weight if mentioned (or null)
- "notes": Any additional info, especially if this is a CHANGE/UPDATE
- "is_update": true if this modifies an existing deadline, false if new

Return ONLY the JSON array, no other text. If no deadlines found, return [].

Example:
[
  {{"date": "2025-10-01", "title": "Assignment 2", "type": "assignment", "weight": 15,
    "notes": "EXTENDED from Sep 24 to Oct 1", "is_update": true}}
]
"""

            try:
                # Use fallback system for automatic model rotation
                response = self._generate_with_fallback(
                    prompt,
                    generation_config=genai.types.GenerationConfig(
                        temperature=0,
                        max_output_tokens=2048,
                    )
                )

                response_text = response.text.strip()

                # Remove markdown code blocks
                if response_text.startswith('```'):
                    response_text = response_text.split('```')[1]
                    if response_text.startswith('json'):
                        response_text = response_text[4:]
                    response_text = response_text.strip()

                updates = json.loads(response_text)

                for update in updates:
                    update['course'] = course_name
                    update['source'] = f"Announcement: {title}"
                    update['announcement_date'] = posted_at

                all_updates.extend(updates)
                logger.info(f"    Found {len(updates)} deadline updates")

            except json.JSONDecodeError:
                logger.debug(f"No valid JSON from announcement: {title}")
            except Exception as e:
                logger.debug(f"Error processing announcement: {e}")

        return all_updates


def parse_assignments_to_deadlines(assignments: List[Dict], course_name: str) -> List[Dict]:
    """Convert Canvas assignments to deadline format"""
    deadlines = []

    for assignment in assignments:
        due_at = assignment.get('due_at')
        if not due_at:
            continue

        # Parse ISO date
        try:
            due_date = datetime.fromisoformat(due_at.replace('Z', '+00:00'))
            date_str = due_date.strftime('%Y-%m-%d')
        except:
            continue

        deadline = {
            'date': date_str,
            'title': assignment.get('name', 'Untitled Assignment'),
            'type': 'assignment',
            'weight': assignment.get('points_possible'),
            'notes': f"Canvas Assignment",
            'course': course_name,
            'source': 'Canvas Assignments'
        }
        deadlines.append(deadline)

    return deadlines


def parse_quizzes_to_deadlines(quizzes: List[Dict], course_name: str) -> List[Dict]:
    """Convert Canvas quizzes to deadline format"""
    deadlines = []

    for quiz in quizzes:
        due_at = quiz.get('due_at')
        if not due_at:
            continue

        # Parse ISO date
        try:
            due_date = datetime.fromisoformat(due_at.replace('Z', '+00:00'))
            date_str = due_date.strftime('%Y-%m-%d')
        except:
            continue

        deadline = {
            'date': date_str,
            'title': quiz.get('title', 'Untitled Quiz'),
            'type': 'quiz',
            'weight': quiz.get('points_possible'),
            'notes': f"Canvas Quiz - {quiz.get('time_limit', 'No')} min time limit" if quiz.get('time_limit') else "Canvas Quiz",
            'course': course_name,
            'source': 'Canvas Quizzes'
        }
        deadlines.append(deadline)

    return deadlines


def parse_calendar_events_to_deadlines(events: List[Dict], course_name: str) -> List[Dict]:
    """Convert Canvas calendar events to deadline format"""
    deadlines = []

    for event in events:
        start_at = event.get('start_at')
        if not start_at:
            continue

        # Parse ISO date
        try:
            event_date = datetime.fromisoformat(start_at.replace('Z', '+00:00'))
            date_str = event_date.strftime('%Y-%m-%d')
        except:
            continue

        # Determine event type
        event_type = 'event'
        title = event.get('title', 'Untitled Event')

        # Try to infer type from title
        title_lower = title.lower()
        if 'exam' in title_lower or 'test' in title_lower:
            event_type = 'exam'
        elif 'assignment' in title_lower or 'hw' in title_lower:
            event_type = 'assignment'
        elif 'quiz' in title_lower:
            event_type = 'quiz'

        deadline = {
            'date': date_str,
            'title': title,
            'type': event_type,
            'weight': None,
            'notes': event.get('description', 'Canvas Calendar Event'),
            'course': course_name,
            'source': 'Canvas Calendar'
        }
        deadlines.append(deadline)

    return deadlines


def get_files_from_modules(modules: List[Dict], canvas_client, course_id: int, max_modules: int = 2) -> List[Dict]:
    """Extract files from first N modules (likely contains intro materials)"""
    intro_files = []

    for i, module in enumerate(modules[:max_modules]):
        module_name = module.get('name', '').lower()

        # Focus on intro/week 0/week 1 modules
        is_intro_module = any(keyword in module_name for keyword in [
            'introduction', 'intro', 'week 0', 'week 1', 'week0', 'week1',
            'getting started', 'syllabus', 'overview', 'orientation'
        ])

        if not is_intro_module and i > 0:
            # Skip modules after first if not explicitly intro
            continue

        items = module.get('items', [])
        for item in items:
            if item.get('type') == 'File':
                # Add file info
                intro_files.append({
                    'id': item.get('content_id'),
                    'filename': item.get('title', ''),
                    'url': item.get('url', ''),
                    'module': module.get('name', '')
                })

    return intro_files


def is_intro_document_ai(filename: str, extractor: 'DeadlineExtractor') -> bool:
    """
    Use AI to classify if a filename is an intro/syllabus document.
    This is smarter than keyword matching and handles edge cases like
    "Course Schedule and Assessments" vs "Quiz 1 Assessment".
    """
    try:
        prompt = f"""Is this filename likely an introductory/syllabus document for a university course?

Filename: {filename}

An intro/syllabus document typically contains:
- Course overview, syllabus, or schedule
- Introduction materials (L0, L01, Week 0, Week 1, Topic 0)
- Module information or course information
- Assessment schedules (list of all assessments for the course)

NOT intro documents:
- Individual tutorials, assignments, or quizzes
- Midterm/final exam papers
- Practice problems or solutions
- Individual lab materials

Answer ONLY with "YES" or "NO"."""

        # Use extractor's fallback system (but with shorter delay since this is just classification)
        time.sleep(0.5)  # Small delay to avoid rate limiting
        response = extractor.model.generate_content(
            prompt,
            generation_config=genai.types.GenerationConfig(
                temperature=0,
                max_output_tokens=10,
            )
        )

        answer = response.text.strip().upper()
        return answer == "YES"

    except Exception as e:
        logger.warning(f"AI classification failed for {filename}, using fallback: {e}")
        # Fallback to simple heuristics
        return is_intro_document_fallback(filename)


def is_intro_document_fallback(filename: str) -> bool:
    """Fallback heuristic-based classification (used when AI fails)"""
    filename_lower = filename.lower()

    # Strong intro indicators always return True
    strong_intro_keywords = [
        'syllabus', 'course outline', 'course schedule',
        'module information', 'course info'
    ]
    if any(keyword in filename_lower for keyword in strong_intro_keywords):
        return True

    # Check regex patterns for L0, L01, Topic 0, etc.
    for pattern in INTRO_PATTERNS:
        if re.search(pattern, filename_lower):
            return True

    # Specific exclusions (but check after strong indicators)
    exclude_keywords = [
        'midterm', 'mid-term', 'mid_term',
        'tutorial', 'tut', 'quiz ',  # Note: space after quiz to avoid excluding "quiz schedule"
        'exam', 'final', 'practice', 'solution', 'answer',
        'homework', 'lab', 'project'
    ]

    if any(exclude in filename_lower for exclude in exclude_keywords):
        return False

    return False


def detect_changes(old_deadlines: List[Dict], new_deadlines: List[Dict]) -> Dict[str, List[Dict]]:
    """Detect changes between old and new deadline lists"""
    changes = {
        'added': [],
        'removed': [],
        'modified': []
    }

    # Create lookup dictionaries (key: course + title)
    def make_key(d):
        return f"{d.get('course', '')}_{d.get('title', '')}".lower()

    old_dict = {make_key(d): d for d in old_deadlines}
    new_dict = {make_key(d): d for d in new_deadlines}

    # Find added and modified
    for key, new_item in new_dict.items():
        if key not in old_dict:
            changes['added'].append(new_item)
        else:
            old_item = old_dict[key]
            # Check if date changed
            if old_item.get('date') != new_item.get('date'):
                change_info = new_item.copy()
                change_info['old_date'] = old_item.get('date')
                changes['modified'].append(change_info)

    # Find removed
    for key, old_item in old_dict.items():
        if key not in new_dict:
            changes['removed'].append(old_item)

    return changes


def format_deadlines_markdown(deadlines: List[Dict], changes: Optional[Dict] = None) -> str:
    """Format deadlines as beautiful markdown with color coding"""
    if not deadlines:
        return "No deadlines found."

    # Sort by date (handle None values)
    sorted_deadlines = sorted(deadlines, key=lambda x: x.get('date') or '9999-12-31')

    # Generate dynamic color palette from actual courses
    unique_courses = list(set(d.get('course', '').split()[0] for d in deadlines if d.get('course')))
    color_palette = ['blue', 'blueviolet', 'orange', 'green', 'red', 'brightgreen',
                     'yellow', 'yellowgreen', 'purple', 'pink', 'cyan', 'teal']
    course_colors = {course: color_palette[i % len(color_palette)]
                     for i, course in enumerate(sorted(unique_courses))}

    md = f"# 📚 NUS Assignment Deadlines\n"
    md += f"**Last Updated:** {datetime.now().strftime('%A, %d %B %Y at %H:%M')}\n\n"
    md += f"**Total Deadlines:** {len([d for d in deadlines if d.get('date')])}\n\n"

    # Add changes summary if provided
    if changes and (changes['added'] or changes['removed'] or changes['modified']):
        md += "## ⚠️ RECENT CHANGES DETECTED\n\n"

        if changes['modified']:
            md += "### 📅 Date Changes:\n"
            for item in changes['modified']:
                old_date = item.get('old_date', 'Unknown')
                new_date = item.get('date', 'Unknown')
                course = item.get('course', 'Unknown')
                title = item.get('title', 'Untitled')
                md += f"- **{course}** - {title}: ~~{old_date}~~ → **{new_date}**\n"
            md += "\n"

        if changes['added']:
            md += "### ✨ New Deadlines:\n"
            for item in changes['added']:
                date = item.get('date', 'TBD')
                course = item.get('course', 'Unknown')
                title = item.get('title', 'Untitled')
                md += f"- **{course}** - {title} ({date})\n"
            md += "\n"

        if changes['removed']:
            md += "### 🗑️ Removed/Cancelled:\n"
            for item in changes['removed']:
                course = item.get('course', 'Unknown')
                title = item.get('title', 'Untitled')
                md += f"- **{course}** - {title}\n"
            md += "\n"

        md += "---\n\n"

    current_month = None
    for deadline in sorted_deadlines:
        date_str = deadline.get('date', 'TBD')

        # Add month header
        if date_str != 'TBD':
            try:
                date_obj = datetime.fromisoformat(date_str)
                month_header = date_obj.strftime('%B %Y')
                if month_header != current_month:
                    current_month = month_header
                    md += f"\n## {month_header}\n\n"

                formatted_date = date_obj.strftime('%a, %d %b')
            except:
                formatted_date = date_str
        else:
            formatted_date = 'TBD'
            if current_month != 'TBD':
                current_month = 'TBD'
                md += f"\n## Date TBD\n\n"

        # Add deadline entry with color coding
        title = deadline.get('title', 'Untitled')
        course = deadline.get('course', 'Unknown Course')
        weight = deadline.get('weight')
        weight_str = f" `{weight}%`" if weight else ""
        deadline_type = deadline.get('type', 'other')
        notes = deadline.get('notes', '')
        is_update = deadline.get('is_update', False)
        source = deadline.get('source', '')

        # Extract course code for color
        course_code = course.split()[0] if course else 'Unknown'
        color = course_colors.get(course_code, 'informational')

        # Type emoji
        type_emoji = {
            'assignment': '📝',
            'quiz': '📊',
            'exam': '📄',
            'project': '🚀',
            'presentation': '🎤',
            'other': '📌'
        }.get(deadline_type.lower(), '📌')

        # Style 3: Minimal (Ultra-Clean)
        md += f"**{formatted_date}** • {course_code} • {title}{weight_str}"

        if notes and (is_update or 'EXTENDED' in notes.upper() or 'CHANGED' in notes.upper()):
            md += f"\n  ⚠️ UPDATE: {notes}"
        elif notes:
            md += f"\n  ℹ️ {notes}"

        md += "\n\n"

    return md


def main():
    """Main execution function"""
    parser = argparse.ArgumentParser(description='NUS Canvas Deadline Scraper')
    parser.add_argument('--display', action='store_true', help='Display deadlines in terminal')
    parser.add_argument('--search', type=str, help='Search keywords (comma-separated)')
    args = parser.parse_args()

    # Validate configuration
    if not CANVAS_API_TOKEN:
        logger.error("CANVAS_API_TOKEN not set in .env file")
        return

    if not GEMINI_API_KEY:
        logger.error("GEMINI_API_KEY not set in .env file")
        return

    logger.info("Starting Canvas Deadline Scraper...")
    logger.info(f"📅 Current semester period: {SEMESTER_START} to {SEMESTER_END}")

    # Clean up temporary files from previous runs
    import shutil
    if DOWNLOAD_DIR.exists():
        shutil.rmtree(DOWNLOAD_DIR)
        DOWNLOAD_DIR.mkdir(exist_ok=True)

    # Load previous deadlines for change detection
    old_deadlines = []
    if Path('deadlines.json').exists():
        try:
            with open('deadlines.json', 'r', encoding='utf-8') as f:
                old_deadlines = json.load(f)
            logger.info(f"Loaded {len(old_deadlines)} previous deadlines for change detection")
        except Exception as e:
            logger.warning(f"Could not load previous deadlines: {e}")

    # Initialize clients
    canvas = CanvasClient(CANVAS_API_URL, CANVAS_API_TOKEN)
    extractor = DeadlineExtractor(GEMINI_API_KEY)

    # Get current user
    user = canvas.get_current_user()
    if not user:
        logger.error("Failed to authenticate with Canvas API")
        return

    logger.info(f"Authenticated as: {user.get('name', 'Unknown')}")

    # Get enrolled courses
    courses = canvas.get_enrolled_courses()
    logger.info(f"Found {len(courses)} enrolled courses")

    all_deadlines = []
    current_year = datetime.now().year

    # Process each course
    for course in courses:
        course_id = course.get('id')
        course_name = course.get('name', f'Course {course_id}')

        logger.info(f"\nProcessing: {course_name}")

        # PRIORITY 1: Canvas Calendar API (most comprehensive source)
        logger.info(f"  [1/7] Fetching calendar events...")
        calendar_events = canvas.get_calendar_events(course_id)
        if calendar_events:
            event_deadlines = parse_calendar_events_to_deadlines(calendar_events, course_name)
            all_deadlines.extend(event_deadlines)
            logger.info(f"  ✓ Found {len(event_deadlines)} deadlines from calendar")

        # PRIORITY 2: Canvas Assignments API (direct source)
        logger.info(f"  [2/7] Fetching assignments from Canvas API...")
        assignments = canvas.get_course_assignments(course_id)
        if assignments:
            assignment_deadlines = parse_assignments_to_deadlines(assignments, course_name)
            all_deadlines.extend(assignment_deadlines)
            logger.info(f"  ✓ Found {len(assignment_deadlines)} assignments with due dates")

        # PRIORITY 3: Canvas Quizzes API (direct source)
        logger.info(f"  [3/7] Fetching quizzes from Canvas API...")
        quizzes = canvas.get_course_quizzes(course_id)
        if quizzes:
            quiz_deadlines = parse_quizzes_to_deadlines(quizzes, course_name)
            all_deadlines.extend(quiz_deadlines)
            logger.info(f"  ✓ Found {len(quiz_deadlines)} quizzes with due dates")

        # PRIORITY 4: Course homepage/syllabus (for deadline info like CS2040C)
        logger.info(f"  [4/7] Checking course homepage/syllabus...")
        course_details = canvas.get_course_details(course_id)
        syllabus_body = course_details.get('syllabus_body', '')
        if syllabus_body and syllabus_body.strip():
            syllabus_text = html_to_text(syllabus_body)
            if syllabus_text.strip():
                time.sleep(6.5)  # Rate limiting
                deadlines = extractor.extract_deadlines(syllabus_text, course_name, current_year)
                if deadlines:
                    all_deadlines.extend(deadlines)
                    logger.info(f"  ✓ Found {len(deadlines)} deadlines from course homepage")

        # PRIORITY 5: Files from first 1-2 Canvas Modules (smartly identified intro materials)
        logger.info(f"  [5/7] Checking first 1-2 course modules for intro materials...")
        modules = canvas.get_course_modules(course_id)
        if modules:
            module_files = get_files_from_modules(modules, canvas, course_id, max_modules=2)
            if module_files:
                logger.info(f"  ✓ Found {len(module_files)} files in intro modules")
                # Process these files (already curated by module position)
                for file_info in module_files[:2]:  # Limit to 2 files max
                    filename = file_info.get('filename', 'unnamed')
                    # For module files, we trust module position, so skip AI classification
                    course_dir = DOWNLOAD_DIR / course_name.replace('/', '_')
                    course_dir.mkdir(exist_ok=True)
                    save_path = course_dir / filename

                    # Get file details and download
                    files = canvas.get_course_files(course_id)
                    file_obj = next((f for f in files if f.get('display_name') == filename), None)
                    if file_obj and file_obj.get('url'):
                        if not save_path.exists():
                            if not canvas.download_file(file_obj['url'], save_path):
                                continue

                        logger.info(f"  Parsing: {filename}")
                        text = DocumentParser.parse_document(save_path)
                        if text.strip():
                            deadlines = extractor.extract_deadlines(text, course_name, current_year)
                            all_deadlines.extend(deadlines)

        # PRIORITY 6: Keyword-based intro PDF detection (fast, no AI needed)
        logger.info(f"  [6/7] Checking for intro PDFs (keyword-based)...")
        files = canvas.get_course_files(course_id)

        if files:
            # Fast keyword matching for intro documents
            intro_keywords = [
                'intro', 'syllabus', 'l0', 'l00', 'l01', 'topic_0', 'topic 0',
                'course schedule', 'course outline', 'assessment', 'module info'
            ]

            intro_files = []
            for f in files[:30]:  # Check first 30 files
                filename = f.get('filename', '').lower()
                if any(keyword in filename for keyword in intro_keywords):
                    intro_files.append(f)
                    if len(intro_files) >= 3:  # Max 3 intro docs
                        break

            if intro_files:
                logger.info(f"  ✓ Found {len(intro_files)} intro document(s) by keyword")

                # Process intro documents
                for file_info in intro_files:
                    filename = file_info.get('filename', 'unnamed')
                    file_url = file_info.get('url')

                    if not file_url:
                        continue

                    # Download file
                    course_dir = DOWNLOAD_DIR / course_name.replace('/', '_')
                    course_dir.mkdir(exist_ok=True)
                    save_path = course_dir / filename

                    if not save_path.exists():
                        if not canvas.download_file(file_url, save_path):
                            continue

                    # Parse document
                    logger.info(f"  Parsing: {filename}")
                    text = DocumentParser.parse_document(save_path)

                    if not text.strip():
                        logger.warning(f"  No text extracted from {filename}")
                        continue

                    # Extract deadlines with retry on safety filter
                    try:
                        deadlines = extractor.extract_deadlines(text, course_name, current_year)
                        all_deadlines.extend(deadlines)
                    except Exception as e:
                        if 'finish_reason' in str(e) and '2' in str(e):
                            logger.warning(f"  Safety filter blocked {filename}, retrying with different model...")
                            # Retry with a different model
                            old_model_index = extractor.current_model_index
                            extractor.current_model_index = (extractor.current_model_index + 1) % len(extractor.available_models)
                            extractor.model = genai.GenerativeModel(extractor.available_models[extractor.current_model_index])
                            try:
                                deadlines = extractor.extract_deadlines(text, course_name, current_year)
                                all_deadlines.extend(deadlines)
                                logger.info(f"  ✓ Retry succeeded with {extractor.available_models[extractor.current_model_index]}")
                            except:
                                logger.error(f"  ✗ Retry failed, skipping {filename}")
                            extractor.current_model_index = old_model_index
                            extractor.model = genai.GenerativeModel(extractor.available_models[extractor.current_model_index])
                        else:
                            logger.error(f"  Failed to extract from {filename}: {e}")

        # PRIORITY 7: Announcements (for last-minute deadline updates)
        logger.info(f"  [7/7] Checking announcements for deadline updates...")
        announcements = canvas.get_course_announcements(course_id, limit=15)
        if announcements:
            updates = extractor.extract_from_announcements(announcements, course_name, current_year)
            if updates:
                logger.info(f"  ✓ Extracted {len(updates)} deadline updates from announcements")
                all_deadlines.extend(updates)

    # Filter deadlines to only include current semester
    logger.info(f"\nFiltering deadlines to semester period: {SEMESTER_START} to {SEMESTER_END}")
    total_before_filter = len(all_deadlines)
    all_deadlines, filtered_out = filter_valid_semester_deadlines(all_deadlines)

    if filtered_out:
        logger.info(f"  ⚠️  Filtered out {len(filtered_out)} deadlines outside semester period")
        logger.info(f"  ✓ Kept {len(all_deadlines)} valid deadlines within semester")

        # Log some examples of filtered deadlines
        for deadline in filtered_out[:5]:
            logger.debug(f"    - {deadline.get('course')}: {deadline.get('title')} ({deadline.get('date')})")
    else:
        logger.info(f"  ✓ All {len(all_deadlines)} deadlines are within semester period")

    # Detect changes
    changes = None
    if old_deadlines:
        logger.info(f"\nDetecting changes from previous run...")
        changes = detect_changes(old_deadlines, all_deadlines)
        total_changes = len(changes['added']) + len(changes['removed']) + len(changes['modified'])
        if total_changes > 0:
            logger.info(f"  🔔 Found {total_changes} changes:")
            logger.info(f"    - {len(changes['added'])} added")
            logger.info(f"    - {len(changes['modified'])} date changes")
            logger.info(f"    - {len(changes['removed'])} removed")
        else:
            logger.info(f"  No changes detected")

    # Save results
    logger.info(f"\nTotal deadlines extracted: {len(all_deadlines)}")

    # Save JSON
    with open('deadlines.json', 'w', encoding='utf-8') as f:
        json.dump(all_deadlines, f, indent=2, ensure_ascii=False)
    logger.info("Saved to: deadlines.json")

    # Save Markdown (with changes highlighted)
    markdown_content = format_deadlines_markdown(all_deadlines, changes)
    with open('deadlines.md', 'w', encoding='utf-8') as f:
        f.write(markdown_content)
    logger.info("Saved to: deadlines.md")

    # Save changes summary if any
    if changes and (changes['added'] or changes['removed'] or changes['modified']):
        with open('changes.json', 'w', encoding='utf-8') as f:
            json.dump(changes, f, indent=2, ensure_ascii=False)
        logger.info("Saved changes to: changes.json")

    # Display if requested
    if args.display:
        print("\n" + "="*80)
        print(markdown_content)
        print("="*80)

    logger.info("\nDone! Check deadlines.md for formatted output.")


if __name__ == "__main__":
    main()
